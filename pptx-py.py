from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt

""" Ref for slide types:  
0 ->  title and subtitle 
1 ->  title and content 
2 ->  section header 
3 ->  two content 
4 ->  Comparison 
5 ->  Title only  
6 ->  Blank 
7 ->  Content with caption 
8 ->  Pic with caption 
"""

# Creating presentation object
root = Presentation()
root.slide_width = Inches(16.675)
root.slide_height = Inches(8.225)
# Creating slide layout
first_slide_layout = root.slide_layouts[0]


slide = root.slides.add_slide(first_slide_layout)

# Adding title and subtitle in
# slide i.e. first page of slide
slide.shapes.title.text = " Created By python-pptx"

# We have different formats of
# subtitles in ppts, for simple
# subtitle this method should
# implemented, you can change
# 0 to 1 for different design
slide.placeholders[1].text = " This is 2nd way"

left = top = Inches(0)

# adding images
pic = slide.shapes.add_picture("./back.png", left, top)

# font
# creating textBox
left = Inches(3)
top = Inches(2)
txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(3))

# creating textFrames
tf = txBox.text_frame
# tf.fit_text(font_family="Montserrat", font_file="./Montserrat-Regular.ttf")
p = tf.add_paragraph()
p.text = "This is a third paragraph that's big "
p.font.size = Pt(20)
p.font.name = "Montserrat"
tf.fit_text(font_family="Montserrat", font_file="./Montserrat-Regular.ttf")
# p.font.bold = True
# p.font.italic = True

# !Background
# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# left = top = Inches(0)
# img_path = declare_your_image_path.jpg
# pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# # This moves it to the background
# slide.shapes._spTree.remove(pic._element)
# slide.shapes._spTree.insert(2, pic._element)

# Saving file
root.save(
    "Output.pptx",
)

print("done")
