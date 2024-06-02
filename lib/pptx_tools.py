from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
import convertapi
from dotenv import dotenv_values
import os
from pptx.enum.shapes import MSO_SHAPE
from lib.utils import get_random_string
from pptx.dml.color import RGBColor

config = dotenv_values(".env")

convertapi.api_secret = config["CONVERT_API_SECRET"]


def convert_pptx_to_image(pptx_path):
    name = pptx_path.split("/")[-1].split(".")[0]
    folder_path = f"./images/{name}"
    os.makedirs(folder_path)

    convertapi.convert("png", {"File": pptx_path}, from_format="pptx").save_files(
        folder_path
    )

    os.rename(f"./images/{name}/{name}.png", f"./images/{name}/{name}-1.png")
    return name


pixels_per_inch = 72


def pixels2inches(pixels):
    return Inches(pixels / pixels_per_inch)


WIDTH = pixels2inches(1200)
HEIGHT = pixels2inches(593)

padding_left = pixels2inches(60)
padding_top = pixels2inches(60)
padding_bottom = pixels2inches(190)


def render_pptx(design):
    root = Presentation()
    root.slide_width = WIDTH
    root.slide_height = HEIGHT

    blank_layout = root.slide_layouts[6]
    slide = root.slides.add_slide(blank_layout)

    # !Background
    slide.shapes.add_picture(design.background.path, pixels2inches(0), pixels2inches(0))

    # !Image 512x512
    slide.shapes.add_picture(
        design.image.path,
        pixels2inches(600 + (600 - 512) / 2),
        pixels2inches((593 - 512) / 2),
    )

    # !Button
    shapes = slide.shapes
    box_heigt_px = 60
    box_width_px = 400
    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        padding_left,
        pixels2inches(593 - 190 - box_heigt_px),
        pixels2inches(box_width_px),
        pixels2inches(box_heigt_px),
    )

    shape.font.fill.solid()
    shape.font.color.rgb = RGBColor(
        design.button.text_color[0],
        design.button.text_color[1],
        design.button.text_color[2],
    )
    shape.text = design.button.text
    shape.font.fill.fore_color.rgb = RGBColor(
        design.button.background_color[0],
        design.button.background_color[1],
        design.button.background_color[2],
    )
    shape.text_frame.line.width = pixels2inches(1)

    # !Header
    header_box = slide.shapes.add_textbox(
        padding_left, padding_top, pixels2inches(600 - 120), Inches(3)
    )
    tf = header_box.text_frame
    p = tf.add_paragraph()
    p.text = design.header.text
    p.font.size = pixels2inches(100)
    p.font.color.rgb = RGBColor(
        design.header.text_color[0],
        design.header.text_color[1],
        design.header.text_color[2],
    )
    # p.font.name = "Montserrat"
    tf.fit_text(font_family="Montserrat", font_file="./Montserrat-Regular.ttf")

    # !Text
    text_box = slide.shapes.add_textbox(
        padding_left, pixels2inches(60 + 100), pixels2inches(600 - 120), Inches(3)
    )
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = design.text.text
    p.font.size = pixels2inches(50)
    p.font.color.rgb = RGBColor(
        design.text.text_color[0], design.text.text_color[1], design.text.text_color[2]
    )
    # p.font.name = "Montserrat"
    tf.fit_text()  # font_family="Montserrat", font_file="./Montserrat-Regular.ttf"

    name = get_random_string(4)
    root.save(
        f"./images/{name}.pptx",
    )
